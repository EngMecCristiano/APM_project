"""
Gerador do Pitch Deck AR² — Asset Reliability & Risks Management
Cria arquivo PPTX com 12 slides profissionais usando python-pptx + matplotlib
"""
import sys
sys.path.insert(0, '/home/cristiano/.local/lib/python3.12/site-packages')

import io
import os
import math
import tempfile
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patches as FancyBboxPatch
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Circle, Arc
from matplotlib.lines import Line2D
import numpy as np
from PIL import Image, ImageDraw, ImageFont

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import pptx.oxml as oxml

# ── Paleta de cores ────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1B, 0x34, 0x64)
NAVY_MED  = RGBColor(0x1E, 0x40, 0x80)
BLUE      = RGBColor(0x2B, 0x5E, 0xAC)
CYAN      = RGBColor(0x00, 0xB4, 0xD8)
GRAY      = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT= RGBColor(0xE5, 0xE7, 0xEB)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG   = RGBColor(0x0D, 0x1B, 0x2A)
GREEN     = RGBColor(0x10, 0xB9, 0x81)
ORANGE    = RGBColor(0xF5, 0x9E, 0x0B)
RED       = RGBColor(0xDC, 0x26, 0x26)
PURPLE    = RGBColor(0x7C, 0x3A, 0xED)

# Hex para matplotlib
C_NAVY    = '#1B3464'
C_NAVYMD  = '#1E4080'
C_BLUE    = '#2B5EAC'
C_CYAN    = '#00B4D8'
C_GRAY    = '#6B7280'
C_DARKBG  = '#0D1B2A'
C_WHITE   = '#FFFFFF'
C_GREEN   = '#10B981'
C_ORANGE  = '#F59E0B'
C_RED     = '#DC2626'
C_PURPLE  = '#7C3AED'
C_GRAYLT  = '#E5E7EB'

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # layout totalmente em branco


# ── Helpers ────────────────────────────────────────────────────────────────────

def fig_to_img(fig):
    """Converte figura matplotlib para BytesIO PNG."""
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_img(slide, img_buf, x, y, w, h=None):
    if h:
        slide.shapes.add_picture(img_buf, x, y, w, h)
    else:
        slide.shapes.add_picture(img_buf, x, y, w)


def bg(slide, color):
    """Preenche o fundo inteiro do slide."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=color)


# ── Logo AR² (recriado com matplotlib) ────────────────────────────────────────

def make_logo_img(dark_bg=True, size=(800, 400)):
    """Recria o logo AR² como imagem PNG."""
    fig, ax = plt.subplots(figsize=(size[0]/100, size[1]/100), dpi=100)
    bg_color = C_DARKBG if dark_bg else C_WHITE
    fig.patch.set_facecolor(bg_color)
    ax.set_facecolor(bg_color)
    ax.set_xlim(0, 8); ax.set_ylim(0, 4)
    ax.axis('off')

    gear_color = '#FFFFFF' if dark_bg else C_NAVY
    text_color = '#FFFFFF' if dark_bg else C_NAVY
    sub_color  = C_CYAN if dark_bg else C_GRAY

    # Gear
    cx, cy, r = 1.6, 2.0, 0.9
    teeth = 10
    for i in range(teeth):
        angle = 2 * math.pi * i / teeth
        x1 = cx + (r) * math.cos(angle)
        y1 = cy + (r) * math.sin(angle)
        x2 = cx + (r + 0.25) * math.cos(angle + math.pi / teeth)
        y2 = cy + (r + 0.25) * math.sin(angle + math.pi / teeth)
        x3 = cx + (r + 0.25) * math.cos(angle + 2 * math.pi / teeth)
        y3 = cy + (r + 0.25) * math.sin(angle + 2 * math.pi / teeth)
        poly = plt.Polygon([[x1,y1],[x2,y2],[x3,y3]], color=gear_color, zorder=2)
        ax.add_patch(poly)
    gear_circle = plt.Circle((cx, cy), r, color=gear_color, zorder=3)
    ax.add_patch(gear_circle)
    inner = plt.Circle((cx, cy), r * 0.45, color=bg_color, zorder=4)
    ax.add_patch(inner)

    # Circuit lines inside gear
    line_color = C_CYAN if dark_bg else C_BLUE
    for i, dy in enumerate([-0.22, 0.0, 0.22]):
        x_start = cx - 0.3
        x_end   = cx + 0.7
        ax.plot([x_start, x_end], [cy + dy, cy + dy],
                color=line_color, lw=1.8, zorder=5)
        ax.plot([x_end], [cy + dy], 'o', color=line_color,
                markersize=5, zorder=6)

    # AR² text
    ax.text(2.9, 1.85, 'AR', fontsize=62, fontweight='bold',
            color=text_color, ha='left', va='bottom',
            fontfamily='DejaVu Sans', zorder=5)
    ax.text(5.35, 2.75, '2', fontsize=32, fontweight='bold',
            color=text_color, ha='left', va='bottom', zorder=5)

    # Separator line
    ax.plot([2.9, 7.8], [1.72, 1.72], color=sub_color, lw=1.5, zorder=5)

    # Subtitle
    ax.text(2.9, 1.35, 'Asset Reliability & Risks Management',
            fontsize=16, color=sub_color, ha='left', va='top',
            fontfamily='DejaVu Sans', zorder=5)

    buf = fig_to_img(fig)
    return buf


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ═══════════════════════════════════════════════════════════════════════════════
def slide_cover():
    slide = prs.slides.add_slide(BLANK)

    # Fundo gradiente simulado com dois retângulos
    bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=CYAN)  # linha topo

    # Logo
    logo_buf = make_logo_img(dark_bg=True, size=(960, 380))
    add_img(slide, logo_buf, Inches(2.0), Inches(0.9), Inches(9.33))

    # Tagline
    add_text(slide, 'AI-Powered Asset Management · ISO 55001 · ISO 14224',
             Inches(1.0), Inches(4.55), Inches(11.33), Inches(0.5),
             size=16, color=CYAN, align=PP_ALIGN.CENTER)

    # Badge
    badge = add_rect(slide, Inches(4.8), Inches(5.25), Inches(3.73), Inches(0.65),
                     fill=BLUE, line=CYAN, line_w=Pt(1))
    add_text(slide, 'Founding Team Pitch  ·  April 2026',
             Inches(4.8), Inches(5.3), Inches(3.73), Inches(0.55),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Linha inferior
    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), fill=NAVY_MED)

    print("  ✓ Slide 1 — Capa")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — O PROBLEMA
# ═══════════════════════════════════════════════════════════════════════════════
def make_problem_img():
    fig, ax = plt.subplots(figsize=(12, 5.5), dpi=120)
    fig.patch.set_facecolor(C_DARKBG)
    ax.set_facecolor(C_DARKBG)
    ax.axis('off')
    ax.set_xlim(0, 12); ax.set_ylim(0, 5.5)

    problems = [
        ('💸', '$50B+/year', 'Custo de paradas não\nplanejadas na América\ndo Norte', C_RED),
        ('⚠️', '70%', 'das PMEs sem sistema\nde confiabilidade ou\ngestão de riscos', C_ORANGE),
        ('📋', 'ISO 55001\nISO 14224', 'Normas complexas,\ninacessíveis para\nequipes sem especialistas', C_CYAN),
        ('👷', '$300–500/h', 'Custo de consultores\nHAZOP/LOPA/RCM\nno mercado atual', C_PURPLE),
    ]

    for i, (icon, stat, desc, color) in enumerate(problems):
        x = 0.3 + i * 2.95

        # Card
        card = FancyBboxPatch((x, 0.3), 2.6, 4.6,
                               boxstyle="round,pad=0.1",
                               facecolor='#162032', edgecolor=color,
                               linewidth=2, zorder=2)
        ax.add_patch(card)

        # Top accent bar
        accent = FancyBboxPatch((x, 4.65), 2.6, 0.25,
                                 boxstyle="round,pad=0.05",
                                 facecolor=color, zorder=3)
        ax.add_patch(accent)

        # Icon
        ax.text(x + 1.3, 4.2, icon, fontsize=26, ha='center', va='center', zorder=4)

        # Stat
        ax.text(x + 1.3, 3.35, stat, fontsize=17, fontweight='bold',
                color=color, ha='center', va='center', zorder=4,
                fontfamily='DejaVu Sans')

        # Divider
        ax.plot([x + 0.3, x + 2.3], [2.95, 2.95], color=color, alpha=0.4, lw=1)

        # Description
        ax.text(x + 1.3, 1.9, desc, fontsize=10.5, color=C_WHITE,
                ha='center', va='center', zorder=4, linespacing=1.5)

    buf = fig_to_img(fig)
    return buf


def slide_problem():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=RED)

    add_text(slide, 'O PROBLEMA', Inches(0.5), Inches(0.15), Inches(12.0), Inches(0.55),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, 'Gestão de ativos industriais é crítica, cara e inacessível para PMEs',
             Inches(0.5), Inches(0.68), Inches(12.0), Inches(0.38),
             size=14, color=CYAN, align=PP_ALIGN.LEFT)

    img = make_problem_img()
    add_img(slide, img, Inches(0.4), Inches(1.15), Inches(12.5))

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), fill=NAVY_MED)
    print("  ✓ Slide 2 — Problema")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — A SOLUÇÃO
# ═══════════════════════════════════════════════════════════════════════════════
def make_solution_img():
    fig, ax = plt.subplots(figsize=(12, 5.0), dpi=120)
    fig.patch.set_facecolor(C_DARKBG)
    ax.set_facecolor(C_DARKBG)
    ax.axis('off')
    ax.set_xlim(0, 12); ax.set_ylim(0, 5.0)

    # Left: Problem side
    ax.text(2.5, 4.6, 'ANTES', fontsize=13, fontweight='bold',
            color=C_RED, ha='center')
    befores = [
        'Planilhas Excel fragmentadas',
        'Consultores externos caros',
        'Sem RUL / predição de falhas',
        'Análises de risco manuais',
        'Dados isolados por setor',
    ]
    for i, txt in enumerate(befores):
        y = 4.0 - i * 0.68
        card = FancyBboxPatch((0.2, y - 0.22), 4.5, 0.52,
                               boxstyle="round,pad=0.05",
                               facecolor='#2A1A1A', edgecolor='#DC2626',
                               linewidth=1.2, zorder=2)
        ax.add_patch(card)
        ax.text(0.55, y + 0.05, '✗', fontsize=12, color=C_RED, va='center')
        ax.text(1.05, y + 0.05, txt, fontsize=10.5, color='#FECACA', va='center')

    # Arrow
    ax.annotate('', xy=(7.0, 2.5), xytext=(5.0, 2.5),
                arrowprops=dict(arrowstyle='->', color=C_CYAN,
                                lw=3, mutation_scale=25))
    ax.text(6.0, 2.85, 'AR²', fontsize=15, fontweight='bold',
            color=C_CYAN, ha='center')

    # Right: After side
    ax.text(9.5, 4.6, 'DEPOIS', fontsize=13, fontweight='bold',
            color=C_GREEN, ha='center')
    afters = [
        'Plataforma SaaS unificada',
        'IA prescritiva 24/7',
        'Predição ML + RUL em tempo real',
        'HAZOP / LOPA / FTA automatizados',
        'Dashboard ISO 55001 integrado',
    ]
    for i, txt in enumerate(afters):
        y = 4.0 - i * 0.68
        card = FancyBboxPatch((7.3, y - 0.22), 4.5, 0.52,
                               boxstyle="round,pad=0.05",
                               facecolor='#0A2A1A', edgecolor='#10B981',
                               linewidth=1.2, zorder=2)
        ax.add_patch(card)
        ax.text(7.65, y + 0.05, '✓', fontsize=12, color=C_GREEN, va='center')
        ax.text(8.15, y + 0.05, txt, fontsize=10.5, color='#A7F3D0', va='center')

    buf = fig_to_img(fig)
    return buf


def slide_solution():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=GREEN)

    add_text(slide, 'A SOLUÇÃO', Inches(0.5), Inches(0.15), Inches(12.0), Inches(0.55),
             size=28, bold=True, color=WHITE)
    add_text(slide, 'Uma plataforma integrada de Confiabilidade + Segurança de Processo com IA',
             Inches(0.5), Inches(0.68), Inches(12.0), Inches(0.38),
             size=14, color=CYAN)

    img = make_solution_img()
    add_img(slide, img, Inches(0.4), Inches(1.1), Inches(12.5))

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), fill=NAVY_MED)
    print("  ✓ Slide 3 — Solução")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — MÓDULOS DO PRODUTO
# ═══════════════════════════════════════════════════════════════════════════════
def make_modules_img():
    fig, ax = plt.subplots(figsize=(12, 5.2), dpi=120)
    fig.patch.set_facecolor(C_DARKBG)
    ax.set_facecolor(C_DARKBG)
    ax.axis('off')
    ax.set_xlim(0, 12); ax.set_ylim(0, 5.2)

    modules = [
        {
            'title': '🔥 Process Safety',
            'color': C_RED,
            'faded': '#2A0F0F',
            'items': [
                ('HAZOP', 'Análise de risco em nós de processo'),
                ('LOPA', 'Camadas de proteção e SIL'),
                ('SIS/SIF', 'Funções instrumentadas de segurança'),
                ('FTA / ETA', 'Árvore de falhas e eventos'),
                ('IA Agent', 'Revisão automática conforme IEC 61511'),
            ]
        },
        {
            'title': '⚙️ Maintenance Engineering',
            'color': C_CYAN,
            'faded': '#0A1A2A',
            'items': [
                ('LDA / Weibull', 'Análise de vida e confiabilidade'),
                ('RAM / RBD', 'Disponibilidade e diagramas de blocos'),
                ('Crow-AMSAA', 'Tendência NHPP em sistemas reparáveis'),
                ('ML Preditivo', 'Random Forest + RUL + Anomalias'),
                ('IA Prescritiva', 'Plano de manutenção via Claude API'),
            ]
        },
    ]

    for i, mod in enumerate(modules):
        x = 0.3 + i * 6.1

        # Module card
        card = FancyBboxPatch((x, 0.2), 5.7, 4.8,
                               boxstyle="round,pad=0.1",
                               facecolor=mod['faded'],
                               edgecolor=mod['color'], linewidth=2.5, zorder=2)
        ax.add_patch(card)

        # Title bar
        title_bar = FancyBboxPatch((x, 4.65), 5.7, 0.35,
                                    boxstyle="round,pad=0.05",
                                    facecolor=mod['color'], zorder=3)
        ax.add_patch(title_bar)
        ax.text(x + 2.85, 4.83, mod['title'], fontsize=13.5, fontweight='bold',
                color=C_WHITE, ha='center', va='center', zorder=4)

        # Items
        for j, (tag, desc) in enumerate(mod['items']):
            y = 4.15 - j * 0.73
            # Tag badge
            badge = FancyBboxPatch((x + 0.3, y - 0.2), 1.3, 0.38,
                                    boxstyle="round,pad=0.04",
                                    facecolor=mod['color'] + '55',
                                    edgecolor=mod['color'], linewidth=1, zorder=4)
            ax.add_patch(badge)
            ax.text(x + 0.95, y - 0.01, tag, fontsize=9.5, fontweight='bold',
                    color=mod['color'], ha='center', va='center', zorder=5)
            ax.text(x + 1.75, y - 0.01, desc, fontsize=10, color='#CBD5E1',
                    ha='left', va='center', zorder=5)

    # Center divider / ISO badge
    badge2 = FancyBboxPatch((5.15, 2.0), 1.7, 1.2,
                              boxstyle="round,pad=0.1",
                              facecolor='#1E3A5F', edgecolor=C_CYAN, linewidth=2, zorder=6)
    ax.add_patch(badge2)
    ax.text(6.0, 2.82, 'ISO', fontsize=11, fontweight='bold',
            color=C_CYAN, ha='center', va='center', zorder=7)
    ax.text(6.0, 2.42, '55001', fontsize=14, fontweight='bold',
            color=C_WHITE, ha='center', va='center', zorder=7)
    ax.text(6.0, 2.08, '14224', fontsize=10, color=C_GRAY,
            ha='center', va='center', zorder=7)

    buf = fig_to_img(fig)
    return buf


def slide_modules():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=BLUE)

    add_text(slide, 'MÓDULOS DO PRODUTO', Inches(0.5), Inches(0.15), Inches(12.0), Inches(0.55),
             size=28, bold=True, color=WHITE)
    add_text(slide, 'Dois módulos complementares · uma plataforma integrada',
             Inches(0.5), Inches(0.68), Inches(12.0), Inches(0.38),
             size=14, color=CYAN)

    img = make_modules_img()
    add_img(slide, img, Inches(0.4), Inches(1.12), Inches(12.5))

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), fill=NAVY_MED)
    print("  ✓ Slide 4 — Módulos")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ARQUITETURA DA PLATAFORMA
# ═══════════════════════════════════════════════════════════════════════════════
def make_arch_img():
    fig, ax = plt.subplots(figsize=(12, 5.2), dpi=120)
    fig.patch.set_facecolor(C_DARKBG)
    ax.set_facecolor(C_DARKBG)
    ax.axis('off')
    ax.set_xlim(0, 12); ax.set_ylim(0, 5.2)

    def box(x, y, w, h, label, sublabel='', color=C_BLUE, text_color=C_WHITE, fontsize=10):
        card = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.08",
                               facecolor=color + '33', edgecolor=color,
                               linewidth=1.8, zorder=3)
        ax.add_patch(card)
        ax.text(x + w/2, y + h/2 + (0.08 if sublabel else 0), label,
                fontsize=fontsize, fontweight='bold', color=text_color,
                ha='center', va='center', zorder=4)
        if sublabel:
            ax.text(x + w/2, y + h/2 - 0.22, sublabel,
                    fontsize=8, color=C_GRAY, ha='center', va='center', zorder=4)

    def arrow(x1, y1, x2, y2):
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle='->', color=C_CYAN,
                                    lw=1.8, mutation_scale=16),
                    zorder=5)

    # Layer labels
    layers = [
        (0.15, 4.4, 'USERS'),
        (0.15, 3.1, 'FRONTEND'),
        (0.15, 1.95, 'BACKEND'),
        (0.15, 0.55, 'DATA / AI'),
    ]
    for x, y, label in layers:
        ax.text(x, y, label, fontsize=7.5, color=C_GRAY, va='center',
                fontweight='bold', rotation=0)
        ax.plot([1.0, 1.0], [0.2, 5.0], color='#2A3A4A', lw=1, zorder=1)

    # Users
    box(1.2, 4.1, 2.0, 0.65, '🏭 Plant Engineer', '', C_GRAY, C_WHITE, 9.5)
    box(3.5, 4.1, 2.0, 0.65, '🔬 Safety Engineer', '', C_GRAY, C_WHITE, 9.5)
    box(5.8, 4.1, 2.0, 0.65, '📊 Reliability Eng.', '', C_GRAY, C_WHITE, 9.5)
    box(8.1, 4.1, 2.2, 0.65, '👔 Plant Manager', '', C_GRAY, C_WHITE, 9.5)

    # Arrows users → frontend
    for x in [2.2, 4.5, 6.8, 9.2]:
        arrow(x, 4.1, x, 3.75)

    # CDN + API Gateway
    box(1.2, 3.05, 3.5, 0.65, 'React + TypeScript', 'Vite · TailwindCSS · Recharts', C_BLUE, C_WHITE)
    box(5.0, 3.05, 2.5, 0.65, 'CDN / Edge', 'Cloudflare', C_GRAY)
    box(7.8, 3.05, 2.5, 0.65, 'API Gateway', 'Auth · Rate Limit', C_PURPLE)

    arrow(3.0, 3.05, 3.0, 2.7)
    arrow(6.05, 3.05, 8.0, 2.7)
    arrow(9.05, 3.05, 9.05, 2.7)

    # Backend services
    box(1.2, 1.9, 2.2, 0.7, 'FastAPI', 'REST + WebSocket', C_GREEN)
    box(3.7, 1.9, 2.2, 0.7, 'ML Engine', 'RF · IsoForest · RUL', C_CYAN)
    box(6.2, 1.9, 2.0, 0.7, 'AI Agent', 'Claude API', C_ORANGE)
    box(8.5, 1.9, 1.7, 0.7, 'Auth', 'JWT · Tenant', C_PURPLE)

    arrow(2.3, 1.9, 2.3, 1.55)
    arrow(4.8, 1.9, 4.8, 1.55)
    arrow(7.2, 1.9, 7.2, 1.55)

    # Data layer
    box(0.8, 0.25, 1.8, 0.75, 'PostgreSQL', 'Multi-tenant RLS', C_BLUE)
    box(2.9, 0.25, 1.8, 0.75, 'Redis', 'Cache · Sessions', C_RED)
    box(5.0, 0.25, 1.8, 0.75, 'MLflow', 'Model Registry', C_GREEN)
    box(7.1, 0.25, 2.0, 0.75, 'Evidently AI', 'Drift Monitor', C_ORANGE)
    box(9.4, 0.25, 2.3, 0.75, 'Railway / AWS', 'Deploy · Scale', C_PURPLE)

    # CI/CD label
    ax.text(11.7, 2.6, 'CI/CD\nGitHub\nActions', fontsize=7.5, color=C_GRAY,
            ha='center', va='center', linespacing=1.4)

    buf = fig_to_img(fig)
    return buf


def slide_architecture():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=CYAN)

    add_text(slide, 'ARQUITETURA DA PLATAFORMA', Inches(0.5), Inches(0.15), Inches(12.0), Inches(0.55),
             size=28, bold=True, color=WHITE)
    add_text(slide, 'Stack moderno · Cloud-native · Multi-tenant · MLOps integrado',
             Inches(0.5), Inches(0.68), Inches(12.0), Inches(0.38),
             size=14, color=CYAN)

    img = make_arch_img()
    add_img(slide, img, Inches(0.3), Inches(1.12), Inches(12.7))

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), fill=NAVY_MED)
    print("  ✓ Slide 5 — Arquitetura")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — FLUXO DO USUÁRIO
# ═══════════════════════════════════════════════════════════════════════════════
def make_flow_img():
    fig, ax = plt.subplots(figsize=(12, 4.8), dpi=120)
    fig.patch.set_facecolor(C_DARKBG)
    ax.set_facecolor(C_DARKBG)
    ax.axis('off')
    ax.set_xlim(0, 12); ax.set_ylim(0, 4.8)

    steps = [
        ('1', '📥 Import\nDados', 'CSV / API\nManual entry', C_BLUE),
        ('2', '⚙️ Feature\nEngineering', 'TBF · TTR\nLags · MAs', C_CYAN),
        ('3', '📊 LDA /\nWeibull', 'β, η, MTTF\nIC 90%', C_GREEN),
        ('4', '🤖 ML\nPrediction', 'RF · RUL\nAnomalias', C_ORANGE),
        ('5', '🔥 Risk\nScore', '0–100\nCrítico/Alto', C_RED),
        ('6', '💡 IA\nPrescritiva', 'Plano de\nações ISO', C_PURPLE),
        ('7', '📋 Relatório\nPDF/Export', 'ISO 14224\nAuditoria', C_GRAY),
    ]

    sx = 0.5
    for i, (num, title, sub, color) in enumerate(steps):
        x = sx + i * 1.65

        # Circle
        circle = Circle((x + 0.55, 3.6), 0.45, color=color, zorder=3)
        ax.add_patch(circle)
        ax.text(x + 0.55, 3.6, num, fontsize=16, fontweight='bold',
                color=C_WHITE, ha='center', va='center', zorder=4)

        # Arrow
        if i < len(steps) - 1:
            ax.annotate('', xy=(x + 1.55, 3.6), xytext=(x + 1.05, 3.6),
                        arrowprops=dict(arrowstyle='->', color=color,
                                        lw=2, mutation_scale=14), zorder=2)

        # Box
        card = FancyBboxPatch((x, 0.4), 1.1, 2.8,
                               boxstyle="round,pad=0.07",
                               facecolor=color + '22', edgecolor=color,
                               linewidth=1.5, zorder=2)
        ax.add_patch(card)

        # Vertical connector
        ax.plot([x + 0.55, x + 0.55], [3.15, 3.2], color=color, lw=2, zorder=2)

        ax.text(x + 0.55, 2.5, title, fontsize=10, fontweight='bold',
                color=color, ha='center', va='center', linespacing=1.4, zorder=4)
        ax.text(x + 0.55, 1.35, sub, fontsize=8.5, color='#94A3B8',
                ha='center', va='center', linespacing=1.5, zorder=4)

    buf = fig_to_img(fig)
    return buf


def slide_flow():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=GREEN)

    add_text(slide, 'FLUXO DO USUÁRIO', Inches(0.5), Inches(0.15), Inches(12.0), Inches(0.55),
             size=28, bold=True, color=WHITE)
    add_text(slide, 'Da importação de dados à prescrição de manutenção em minutos',
             Inches(0.5), Inches(0.68), Inches(12.0), Inches(0.38),
             size=14, color=CYAN)

    img = make_flow_img()
    add_img(slide, img, Inches(0.3), Inches(1.12), Inches(12.7))

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), fill=NAVY_MED)
    print("  ✓ Slide 6 — Fluxo")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MLOPS PIPELINE
# ═══════════════════════════════════════════════════════════════════════════════
def make_mlops_img():
    fig, ax = plt.subplots(figsize=(12, 4.8), dpi=120)
    fig.patch.set_facecolor(C_DARKBG)
    ax.set_facecolor(C_DARKBG)
    ax.axis('off')
    ax.set_xlim(0, 12); ax.set_ylim(0, 4.8)

    def mlbox(x, y, w, h, title, subs, color):
        card = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.1",
                               facecolor=color + '25', edgecolor=color,
                               linewidth=2, zorder=3)
        ax.add_patch(card)
        ax.text(x + w/2, y + h - 0.22, title, fontsize=10.5, fontweight='bold',
                color=color, ha='center', va='center', zorder=4)
        ax.plot([x + 0.15, x + w - 0.15], [y + h - 0.42, y + h - 0.42],
                color=color, alpha=0.4, lw=1)
        for i, s in enumerate(subs):
            ax.text(x + w/2, y + h - 0.68 - i * 0.35, s, fontsize=8.5,
                    color='#CBD5E1', ha='center', va='center', zorder=4)

    def marrow(x1, x2, y=2.4):
        ax.annotate('', xy=(x2, y), xytext=(x1, y),
                    arrowprops=dict(arrowstyle='->', color=C_CYAN,
                                    lw=2, mutation_scale=18), zorder=5)

    # Pipeline stages
    mlbox(0.2, 0.8, 2.0, 3.2, 'DATA\nINGESTION',
          ['CSV Upload', 'API Push', 'Validation', 'S3 / DB'], C_BLUE)
    marrow(2.2, 2.5)
    mlbox(2.5, 0.8, 2.0, 3.2, 'FEATURE\nSTORE',
          ['TBF/TTR', 'Rolling stats', 'Lag features', 'Redis cache'], C_CYAN)
    marrow(4.5, 4.8)
    mlbox(4.8, 0.8, 2.0, 3.2, 'TRAINING',
          ['RandomForest', 'IsoForest', 'Time-split 80/20', 'MLflow log'], C_GREEN)
    marrow(6.8, 7.1)
    mlbox(7.1, 0.8, 2.0, 3.2, 'SERVING',
          ['RUL Predict', 'Risk Score', 'Anomaly flag', 'REST API'], C_ORANGE)
    marrow(9.1, 9.4)
    mlbox(9.4, 0.8, 2.3, 3.2, 'MONITORING',
          ['Evidently AI', 'PSI drift', 'Auto-retrain', 'Alerts'], C_RED)

    # CI/CD feedback loop
    ax.annotate('', xy=(0.6, 0.8), xytext=(11.2, 0.8),
                arrowprops=dict(arrowstyle='->', color='#374151',
                                lw=1.5, mutation_scale=12,
                                connectionstyle='arc3,rad=0.4'), zorder=2)
    ax.text(6.0, 0.2, 'Drift detected → auto-retrain loop', fontsize=8.5,
            color=C_GRAY, ha='center')

    # Multi-tenant label
    ax.text(0.3, 4.55, 'Cold-start:  < 5 falhas → modelo global    5–20 → fine-tune    > 20 → modelo específico por ativo',
            fontsize=8.5, color=C_GRAY, va='center')

    buf = fig_to_img(fig)
    return buf


def slide_mlops():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=ORANGE)

    add_text(slide, 'MLOPS PIPELINE', Inches(0.5), Inches(0.15), Inches(12.0), Inches(0.55),
             size=28, bold=True, color=WHITE)
    add_text(slide, 'Ciclo completo: ingestão → treinamento → serving → monitoramento de drift',
             Inches(0.5), Inches(0.68), Inches(12.0), Inches(0.38),
             size=14, color=CYAN)

    img = make_mlops_img()
    add_img(slide, img, Inches(0.3), Inches(1.12), Inches(12.7))

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), fill=NAVY_MED)
    print("  ✓ Slide 7 — MLOps")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — MERCADO
# ═══════════════════════════════════════════════════════════════════════════════
def make_market_img():
    fig, axes = plt.subplots(1, 2, figsize=(12, 4.6), dpi=120)
    fig.patch.set_facecolor(C_DARKBG)

    # Left: TAM/SAM/SOM circles
    ax = axes[0]
    ax.set_facecolor(C_DARKBG)
    ax.set_xlim(-4.5, 4.5); ax.set_ylim(-4.5, 4.5)
    ax.axis('off')
    ax.set_aspect('equal')

    circles = [
        (4.0, C_NAVY, '$8.2B', 'TAM\nGlobal Asset\nManagement Software'),
        (2.7, C_BLUE, '$1.4B', 'SAM\nBrazil + Canada\nSME Industrial'),
        (1.4, C_CYAN, '$48M', 'SOM\nTarget 5 years\nPMEs Tier 1'),
    ]
    for r, c, val, label in circles:
        circle = Circle((0, 0), r, facecolor=c + '40', edgecolor=c, linewidth=2)
        ax.add_patch(circle)
    # Labels
    ax.text(0, 0.1, '$48M', fontsize=16, fontweight='bold', color=C_CYAN,
            ha='center', va='center')
    ax.text(0, -0.45, 'SOM', fontsize=9, color=C_CYAN, ha='center')
    ax.text(0, 2.0, '$1.4B  SAM', fontsize=10, color=C_WHITE, ha='center')
    ax.text(0, 3.3, '$8.2B  TAM', fontsize=10, color=C_WHITE, ha='center')

    ax.text(0, -4.2, 'CAGR 12.4% · Source: MarketsandMarkets 2025',
            fontsize=7.5, color=C_GRAY, ha='center')

    # Right: Geography bar
    ax2 = axes[1]
    ax2.set_facecolor(C_DARKBG)
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)
    ax2.spines['left'].set_color('#374151')
    ax2.spines['bottom'].set_color('#374151')
    ax2.tick_params(colors=C_GRAY, labelsize=9)

    segments = ['O&G\nBrasil', 'Mining\nBrasil', 'O&G\nCanadá', 'Chemical\nCanadá', 'Power\nBrasil']
    values   = [18, 12, 9, 7, 5]
    colors   = [C_CYAN, C_BLUE, C_GREEN, C_ORANGE, C_PURPLE]
    bars = ax2.barh(segments, values, color=colors, edgecolor='none', height=0.55)
    for bar, val in zip(bars, values):
        ax2.text(val + 0.3, bar.get_y() + bar.get_height()/2,
                 f'${val}M', color=C_WHITE, va='center', fontsize=9, fontweight='bold')
    ax2.set_xlim(0, 25)
    ax2.set_xlabel('Mercado Endereçável (USD M)', color=C_GRAY, fontsize=8.5)
    ax2.set_title('Segmentos Prioritários', color=C_WHITE, fontsize=11, fontweight='bold', pad=8)
    ax2.set_facecolor(C_DARKBG)
    fig.patch.set_facecolor(C_DARKBG)
    plt.tight_layout(pad=1.5)

    buf = fig_to_img(fig)
    return buf


def slide_market():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=GREEN)

    add_text(slide, 'OPORTUNIDADE DE MERCADO', Inches(0.5), Inches(0.15), Inches(12.0), Inches(0.55),
             size=28, bold=True, color=WHITE)
    add_text(slide, 'Foco: PMEs industriais no Brasil e Canadá — mercado sub-atendido e em crescimento',
             Inches(0.5), Inches(0.68), Inches(12.0), Inches(0.38),
             size=14, color=CYAN)

    img = make_market_img()
    add_img(slide, img, Inches(0.3), Inches(1.12), Inches(12.7))

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), fill=NAVY_MED)
    print("  ✓ Slide 8 — Mercado")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MODELO DE NEGÓCIO
# ═══════════════════════════════════════════════════════════════════════════════
def make_pricing_img():
    fig, ax = plt.subplots(figsize=(12, 4.8), dpi=120)
    fig.patch.set_facecolor(C_DARKBG)
    ax.set_facecolor(C_DARKBG)
    ax.axis('off')
    ax.set_xlim(0, 12); ax.set_ylim(0, 4.8)

    plans = [
        {
            'name': 'STARTER', 'price': '$149', 'period': '/month',
            'color': C_BLUE, 'badge': '',
            'features': [
                '✓ Maintenance Engineering',
                '✓ LDA / Weibull / RAM',
                '✓ ML Predictive (5 assets)',
                '✓ Expert System (no AI)',
                '✓ CSV Export',
                '✗ Process Safety',
                '✗ AI Agent (Claude)',
            ],
        },
        {
            'name': 'PROFESSIONAL', 'price': '$399', 'period': '/month',
            'color': C_CYAN, 'badge': '★ POPULAR',
            'features': [
                '✓ Everything in Starter',
                '✓ Process Safety (HAZOP/LOPA)',
                '✓ Unlimited assets',
                '✓ AI Agent (Claude API)',
                '✓ ISO 14224 Audit Export',
                '✓ Priority Support',
                '✗ Custom Integrations',
            ],
        },
        {
            'name': 'ENTERPRISE', 'price': 'Custom', 'period': '',
            'color': C_PURPLE, 'badge': '',
            'features': [
                '✓ Everything in Professional',
                '✓ White-label',
                '✓ API integrations (SAP/OSI)',
                '✓ On-premise / VPC deploy',
                '✓ Dedicated CSM',
                '✓ SLA 99.9%',
                '✓ Custom ML models',
            ],
        },
    ]

    for i, plan in enumerate(plans):
        x = 0.4 + i * 3.85

        is_popular = plan['badge'] == '★ POPULAR'
        edge_w = 3.0 if is_popular else 1.8

        card = FancyBboxPatch((x, 0.2), 3.4, 4.4,
                               boxstyle="round,pad=0.1",
                               facecolor=plan['color'] + ('30' if is_popular else '18'),
                               edgecolor=plan['color'], linewidth=edge_w, zorder=2)
        ax.add_patch(card)

        # Popular badge
        if is_popular:
            badge = FancyBboxPatch((x + 0.8, 4.35), 1.8, 0.35,
                                    boxstyle="round,pad=0.05",
                                    facecolor=plan['color'], zorder=4)
            ax.add_patch(badge)
            ax.text(x + 1.7, 4.525, plan['badge'], fontsize=8.5, fontweight='bold',
                    color=C_WHITE, ha='center', va='center', zorder=5)

        # Plan name
        ax.text(x + 1.7, 3.9, plan['name'], fontsize=11, fontweight='bold',
                color=plan['color'], ha='center', va='center', zorder=3)

        # Price
        ax.text(x + 1.7, 3.35, plan['price'], fontsize=22, fontweight='bold',
                color=C_WHITE, ha='center', va='center', zorder=3)
        if plan['period']:
            ax.text(x + 1.7, 2.92, plan['period'], fontsize=9, color=C_GRAY,
                    ha='center', va='center', zorder=3)

        # Divider
        ax.plot([x + 0.2, x + 3.2], [2.7, 2.7], color=plan['color'], alpha=0.4, lw=1)

        # Features
        for j, feat in enumerate(plan['features']):
            y = 2.45 - j * 0.33
            color = '#A7F3D0' if feat.startswith('✓') else '#6B7280'
            ax.text(x + 0.25, y, feat, fontsize=8.5, color=color, va='center', zorder=3)

    buf = fig_to_img(fig)
    return buf


def slide_business():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=PURPLE)

    add_text(slide, 'MODELO DE NEGÓCIO', Inches(0.5), Inches(0.15), Inches(12.0), Inches(0.55),
             size=28, bold=True, color=WHITE)
    add_text(slide, 'SaaS B2B · Receita recorrente mensal · 3 tiers de acesso',
             Inches(0.5), Inches(0.68), Inches(12.0), Inches(0.38),
             size=14, color=CYAN)

    img = make_pricing_img()
    add_img(slide, img, Inches(0.3), Inches(1.12), Inches(12.7))

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), fill=NAVY_MED)
    print("  ✓ Slide 9 — Modelo de Negócio")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ROADMAP
# ═══════════════════════════════════════════════════════════════════════════════
def make_roadmap_img():
    fig, ax = plt.subplots(figsize=(12, 4.8), dpi=120)
    fig.patch.set_facecolor(C_DARKBG)
    ax.set_facecolor(C_DARKBG)
    ax.axis('off')
    ax.set_xlim(0, 12); ax.set_ylim(0, 4.8)

    # Timeline line
    ax.plot([0.5, 11.5], [2.8, 2.8], color='#374151', lw=3, zorder=1)

    phases = [
        {
            'x': 2.0, 'label': 'FASE 1', 'period': 'Q2–Q4 2026',
            'title': 'MVP & Validação',
            'color': C_CYAN,
            'items': ['React SPA + FastAPI', 'Módulo Maintenance', 'IA Prescritiva (Claude)',
                      'Deploy Railway', '5–10 clientes beta', 'MRR alvo: $5K'],
            'infra': '$150–300/mês',
        },
        {
            'x': 6.0, 'label': 'FASE 2', 'period': '2027',
            'title': 'Escala',
            'color': C_GREEN,
            'items': ['Process Safety module', 'Multi-tenant robusto', 'Celery + Feast',
                      'Certificação ISO', '50–100 clientes', 'MRR alvo: $40K'],
            'infra': '$800–2.5K/mês',
        },
        {
            'x': 10.0, 'label': 'FASE 3', 'period': '2028+',
            'title': 'Enterprise',
            'color': C_PURPLE,
            'items': ['AWS EKS + SageMaker', 'White-label / API', 'SAP / OSI integração',
                      'Expansão LatAm', '500+ clientes', 'MRR alvo: $300K'],
            'infra': '$3K–15K/mês',
        },
    ]

    for ph in phases:
        x = ph['x']
        color = ph['color']

        # Milestone dot
        dot = Circle((x, 2.8), 0.18, color=color, zorder=4)
        ax.add_patch(dot)

        # Vertical connector
        ax.plot([x, x], [2.98, 3.3], color=color, lw=2, zorder=3)

        # Card above
        card = FancyBboxPatch((x - 1.7, 3.3), 3.4, 1.35,
                               boxstyle="round,pad=0.1",
                               facecolor=color + '22', edgecolor=color,
                               linewidth=2, zorder=3)
        ax.add_patch(card)

        ax.text(x, 4.5, ph['label'], fontsize=10, fontweight='bold',
                color=color, ha='center', va='center', zorder=4)
        ax.text(x, 4.15, ph['period'], fontsize=8.5, color=C_WHITE,
                ha='center', va='center', zorder=4)
        ax.text(x, 3.8, ph['title'], fontsize=9, color=C_GRAY,
                ha='center', va='center', zorder=4)

        # Infra cost
        ax.text(x, 3.42, f'Infra: {ph["infra"]}', fontsize=7.5, color=C_GRAY,
                ha='center', va='center', zorder=4)

        # Items below
        ax.plot([x, x], [2.62, 2.3], color=color, lw=2, zorder=3)
        card2 = FancyBboxPatch((x - 1.7, 0.15), 3.4, 2.15,
                                boxstyle="round,pad=0.1",
                                facecolor='#162032', edgecolor=color + '88',
                                linewidth=1.2, zorder=3)
        ax.add_patch(card2)
        for j, item in enumerate(ph['items']):
            y = 2.1 - j * 0.31
            ax.text(x, y, f'· {item}', fontsize=8, color='#CBD5E1',
                    ha='center', va='center', zorder=4)

    buf = fig_to_img(fig)
    return buf


def slide_roadmap():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=CYAN)

    add_text(slide, 'ROADMAP', Inches(0.5), Inches(0.15), Inches(12.0), Inches(0.55),
             size=28, bold=True, color=WHITE)
    add_text(slide, 'Da validação ao Enterprise em 3 fases — crescimento sustentável',
             Inches(0.5), Inches(0.68), Inches(12.0), Inches(0.38),
             size=14, color=CYAN)

    img = make_roadmap_img()
    add_img(slide, img, Inches(0.3), Inches(1.12), Inches(12.7))

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), fill=NAVY_MED)
    print("  ✓ Slide 10 — Roadmap")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — TIME
# ═══════════════════════════════════════════════════════════════════════════════
def make_team_img():
    fig, ax = plt.subplots(figsize=(12, 4.6), dpi=120)
    fig.patch.set_facecolor(C_DARKBG)
    ax.set_facecolor(C_DARKBG)
    ax.axis('off')
    ax.set_xlim(0, 12); ax.set_ylim(0, 4.6)

    team = [
        {
            'initials': 'CR',
            'role': 'Co-Founder & CEO',
            'title': 'Reliability &\nSafety Domain Expert',
            'color': C_CYAN,
            'skills': ['ISO 14224 / ISO 55001', 'LDA · Weibull · RAM · RCM',
                        'HAZOP / LOPA / SIS', 'APM Analytics (MVP)'],
            'location': '🇧🇷 Brasil → Toronto',
        },
        {
            'initials': 'SR',
            'role': 'Co-Founder & CTO',
            'title': 'Senior Software\nEngineer · Toronto',
            'color': C_GREEN,
            'skills': ['React · TypeScript · Node', 'FastAPI · Microservices',
                        'DevOps · CI/CD · Docker', 'Cloud Architecture (AWS/GCP)'],
            'location': '🇨🇦 Toronto',
        },
        {
            'initials': 'FM',
            'role': 'Co-Founder & CFO',
            'title': 'Finance & Data\nAnalytics · U of T',
            'color': C_ORANGE,
            'skills': ['Financial Modeling / SaaS', 'Data Analytics (MSc U of T)',
                        'Investor Relations', 'Market Strategy BR + CA'],
            'location': '🇨🇦 Toronto',
        },
    ]

    for i, member in enumerate(team):
        x = 0.5 + i * 3.8

        # Card
        card = FancyBboxPatch((x, 0.15), 3.4, 4.25,
                               boxstyle="round,pad=0.1",
                               facecolor=member['color'] + '18',
                               edgecolor=member['color'], linewidth=2, zorder=2)
        ax.add_patch(card)

        # Avatar circle
        avatar = Circle((x + 1.7, 3.7), 0.55, facecolor=member['color'] + '44',
                         edgecolor=member['color'], linewidth=2.5, zorder=3)
        ax.add_patch(avatar)
        ax.text(x + 1.7, 3.7, member['initials'], fontsize=18, fontweight='bold',
                color=member['color'], ha='center', va='center', zorder=4)

        # Role
        ax.text(x + 1.7, 2.98, member['role'], fontsize=9.5, fontweight='bold',
                color=member['color'], ha='center', va='center', zorder=3)

        # Title
        ax.text(x + 1.7, 2.55, member['title'], fontsize=9, color=C_WHITE,
                ha='center', va='center', linespacing=1.4, zorder=3)

        # Divider
        ax.plot([x + 0.3, x + 3.1], [2.25, 2.25], color=member['color'], alpha=0.35, lw=1)

        # Skills
        for j, skill in enumerate(member['skills']):
            y = 2.0 - j * 0.38
            ax.text(x + 0.5, y, f'· {skill}', fontsize=8.5, color='#94A3B8',
                    va='center', zorder=3)

        # Location
        ax.text(x + 1.7, 0.32, member['location'], fontsize=8.5, color=C_GRAY,
                ha='center', va='center', zorder=3)

    buf = fig_to_img(fig)
    return buf


def slide_team():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=ORANGE)

    add_text(slide, 'O TIME', Inches(0.5), Inches(0.15), Inches(12.0), Inches(0.55),
             size=28, bold=True, color=WHITE)
    add_text(slide, 'Domain expertise + Engineering + Finance · Três co-fundadores complementares',
             Inches(0.5), Inches(0.68), Inches(12.0), Inches(0.38),
             size=14, color=CYAN)

    img = make_team_img()
    add_img(slide, img, Inches(0.3), Inches(1.12), Inches(12.7))

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), fill=NAVY_MED)
    print("  ✓ Slide 11 — Time")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — PRÓXIMOS PASSOS
# ═══════════════════════════════════════════════════════════════════════════════
def slide_next():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=CYAN)

    # Logo mini (light version)
    logo_buf = make_logo_img(dark_bg=True, size=(560, 200))
    add_img(slide, logo_buf, Inches(4.5), Inches(0.1), Inches(4.33))

    add_text(slide, 'PRÓXIMOS PASSOS', Inches(0.5), Inches(1.65), Inches(12.0), Inches(0.6),
             size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    steps_data = [
        ('01', 'Formalizar empresa em Toronto (Canada Corp.)', CYAN),
        ('02', 'Refinar MVP APM Analytics — base tecnologica do produto', GREEN),
        ('03', 'Definir stack React + FastAPI — iniciar Sprint 1', BLUE),
        ('04', 'Buscar primeiros 5 clientes beta (industria BR/CA)', ORANGE),
        ('05', 'Pitch para aceleradoras: MaRS Innovation · IRAP · SENAI', PURPLE),
    ]

    for i, (icon, text, color) in enumerate(steps_data):
        x_pos = Inches(2.0)
        y_pos = Inches(2.45 + i * 0.85)

        add_rect(slide, x_pos, y_pos, Inches(0.55), Inches(0.55),
                 fill=color, line=color, line_w=Pt(0))
        add_text(slide, icon, x_pos, y_pos + Inches(0.04), Inches(0.55), Inches(0.48),
                 size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_text(slide, text, x_pos + Inches(0.7), y_pos + Inches(0.04),
                 Inches(8.5), Inches(0.48), size=15, color=WHITE, bold=(i==0))

    add_text(slide, 'Reunião de fundadores · Abril 2026 · Toronto',
             Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.35),
             size=11, color=GRAY, align=PP_ALIGN.CENTER)

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), fill=NAVY_MED)
    print("  ✓ Slide 12 — Próximos Passos")


# ═══════════════════════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════════════════════
print("Gerando slides AR² Pitch Deck...")
slide_cover()
slide_problem()
slide_solution()
slide_modules()
slide_architecture()
slide_flow()
slide_mlops()
slide_market()
slide_business()
slide_roadmap()
slide_team()
slide_next()

output_path = '/home/cristiano/01- Projetos_MLE/APM_project/AR2_PitchDeck_April2026.pptx'
prs.save(output_path)
print(f"\n✅ Salvo em: {output_path}")
print(f"   Slides: {len(prs.slides)}")
